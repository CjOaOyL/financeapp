"""
Build a PowerPoint presentation about Jaquan K. Levons' digital experience
for a mentor meet-and-greet.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Colour palette ──────────────────────────────────────────────────────────
NAVY      = RGBColor(0x1B, 0x3A, 0x5C)
DARK_NAVY = RGBColor(0x0F, 0x24, 0x3E)
ACCENT    = RGBColor(0x2E, 0x86, 0xC1)  # bright blue
GOLD      = RGBColor(0xD4, 0xA0, 0x17)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG  = RGBColor(0xF0, 0xF4, 0xF8)
DARK_TEXT  = RGBColor(0x2C, 0x3E, 0x50)
MED_GRAY  = RGBColor(0x7F, 0x8C, 0x8D)
LIGHT_GRAY = RGBColor(0xBD, 0xC3, 0xC7)
GREEN     = RGBColor(0x27, 0xAE, 0x60)
ORANGE    = RGBColor(0xE6, 0x7E, 0x22)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height

# ── Helper functions ────────────────────────────────────────────────────────

def add_solid_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                font_name='Calibri', line_spacing=1.15):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(0)
    p.space_before = Pt(0)
    if line_spacing != 1.15:
        p.line_spacing = Pt(font_size * line_spacing)
    return txBox

def add_bullet_frame(slide, left, top, width, height, items,
                     font_size=16, color=WHITE, font_name='Calibri',
                     bullet_char='\u25B8', spacing=6, bold_items=False):
    """Add a textbox with bullet items."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = '{} {}'.format(bullet_char, item)
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.font.bold = bold_items
        p.space_after = Pt(spacing)
        p.space_before = Pt(0)
    return txBox

def add_rect(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_rounded_rect(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape

def card_with_header(slide, left, top, width, height, header_text, body_items,
                     header_color=ACCENT, body_color=WHITE, text_color=WHITE,
                     header_font=14, body_font=13, body_text_color=DARK_TEXT):
    """A rounded card with a colored header bar and bullet body."""
    # Card background
    card = add_rounded_rect(slide, left, top, width, height, body_color)
    # Header bar
    add_rect(slide, left, top, width, Inches(0.45), header_color)
    # Header text
    add_textbox(slide, left + Inches(0.15), top + Inches(0.05),
                width - Inches(0.3), Inches(0.4), header_text,
                font_size=header_font, color=WHITE, bold=True,
                alignment=PP_ALIGN.LEFT)
    # Body bullets
    add_bullet_frame(slide, left + Inches(0.15), top + Inches(0.5),
                     width - Inches(0.3), height - Inches(0.6),
                     body_items, font_size=body_font, color=body_text_color,
                     bullet_char='\u2022', spacing=4)
    return card


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_solid_bg(slide, DARK_NAVY)

# Accent line
add_rect(slide, Inches(1.2), Inches(2.2), Inches(1.5), Inches(0.06), GOLD)

add_textbox(slide, Inches(1.2), Inches(2.4), Inches(10), Inches(1.2),
            'JAQUAN K. LEVONS', font_size=48, color=WHITE, bold=True)

add_textbox(slide, Inches(1.2), Inches(3.5), Inches(10), Inches(0.6),
            'Scientific Director  |  Automation, Digital & AI Transformation',
            font_size=24, color=ACCENT, bold=False)

add_textbox(slide, Inches(1.2), Inches(4.2), Inches(10), Inches(0.5),
            'Bristol Myers Squibb  \u2022  Drug Product Development',
            font_size=20, color=LIGHT_GRAY)

# Bottom tagline
add_textbox(slide, Inches(1.2), Inches(5.5), Inches(10), Inches(0.5),
            'My Digital Journey  \u2014  A Mentor Meet & Greet',
            font_size=18, color=GOLD, bold=False)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 2 — ABOUT ME AT A GLANCE
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, LIGHT_BG)

add_textbox(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.6),
            'About Me at a Glance', font_size=36, color=NAVY, bold=True)
add_rect(slide, Inches(0.8), Inches(0.9), Inches(2.0), Inches(0.05), GOLD)

# Left column — bio
bio_items = [
    'B.S.E., Chemical Engineering \u2014 Princeton University (2003)',
    '22+ years at Bristol Myers Squibb in Drug Product Development',
    'Progressed from bench scientist through automation pioneer to digital leader',
    'Currently: Scientific Director, Automation, Digital & AI Transformation',
    'Lead a team of 7 direct reports (3 PhDs) + 3 people managers',
    'Spanning automation, data science, PAT, and predictive modeling',
]
add_bullet_frame(slide, Inches(0.8), Inches(1.2), Inches(5.5), Inches(3.0),
                 bio_items, font_size=17, color=DARK_TEXT, bullet_char='\u25B8',
                 spacing=8)

# Right column — key stats in cards
stat_data = [
    ('22+', 'Years at BMS'),
    ('80+', 'Assets Impacted'),
    ('13', 'Automation Workflows'),
    ('40+', 'GenAI Use Cases Led'),
]
for i, (num, label) in enumerate(stat_data):
    row = i // 2
    col = i % 2
    x = Inches(7.2) + col * Inches(2.8)
    y = Inches(1.3) + row * Inches(2.0)
    box = add_rounded_rect(slide, x, y, Inches(2.5), Inches(1.6), WHITE, ACCENT)
    add_textbox(slide, x, y + Inches(0.15), Inches(2.5), Inches(0.8),
                num, font_size=40, color=ACCENT, bold=True,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x, y + Inches(0.9), Inches(2.5), Inches(0.5),
                label, font_size=14, color=MED_GRAY, bold=False,
                alignment=PP_ALIGN.CENTER)

# Bottom quote
add_textbox(slide, Inches(0.8), Inches(5.8), Inches(11), Inches(0.8),
            '\u201cJaquan is a visionary \u2013 this is undisputed in his feedback.\u201d  \u2014  Senior Leadership',
            font_size=16, color=MED_GRAY, bold=False, alignment=PP_ALIGN.LEFT)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 3 — DIGITAL CAREER JOURNEY  (Timeline)
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, DARK_NAVY)

add_textbox(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.6),
            'My Digital Career Journey', font_size=36, color=WHITE, bold=True)
add_rect(slide, Inches(0.8), Inches(0.9), Inches(2.0), Inches(0.05), GOLD)

# Timeline bar
bar_y = Inches(2.8)
add_rect(slide, Inches(0.5), bar_y, Inches(12.3), Inches(0.05), ACCENT)

# Timeline entries
timeline = [
    ('2003\u20132013', 'Bench Scientist &\nAutomation Developer',
     'Built foundational automation\nfrom scratch: EC, viscous solubility,\norganic-solvent instrument'),
    ('2013\u20132019', 'Automation Leader\n(SRS I \u2192 SRS II)',
     'Grew function to 13 workflows,\n$3M biologics buildout,\nThree-Tier Value Framework'),
    ('2019\u20132022', 'Digital Strategist\n(Principal Scientist)',
     'Created Digital Engineering\ndiscipline, GPS Digital Moonshot,\nARGILE AR/VR (17 sites)'),
    ('2022\u2013Now', 'Digital & AI Leader\n(Scientific Director)',
     'GenAI strategy, Bio In Silico\nmodeling, SPARTAN, C3.AI,\nDevEx Modeling portfolio'),
]

for i, (period, title, details) in enumerate(timeline):
    x = Inches(0.6) + i * Inches(3.1)
    # Circle node
    node = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(1.15),
                                   bar_y - Inches(0.12), Inches(0.25), Inches(0.25))
    node.fill.solid()
    node.fill.fore_color.rgb = GOLD
    node.line.fill.background()
    # Period above
    add_textbox(slide, x, bar_y - Inches(0.65), Inches(2.8), Inches(0.4),
                period, font_size=14, color=GOLD, bold=True, alignment=PP_ALIGN.CENTER)
    # Title below
    add_textbox(slide, x, bar_y + Inches(0.3), Inches(2.8), Inches(0.7),
                title, font_size=15, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    # Details below title
    add_textbox(slide, x, bar_y + Inches(1.1), Inches(2.8), Inches(1.2),
                details, font_size=12, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 4 — THREE-TIER AUTOMATION VALUE FRAMEWORK
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, LIGHT_BG)

add_textbox(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.6),
            'Three-Tier Automation Value Framework', font_size=36, color=NAVY, bold=True)
add_rect(slide, Inches(0.8), Inches(0.9), Inches(2.0), Inches(0.05), GOLD)
add_textbox(slide, Inches(0.8), Inches(1.0), Inches(10), Inches(0.4),
            'An original strategic framework I developed to guide automation investment decisions',
            font_size=15, color=MED_GRAY)

tier_data = [
    ('TIER 1', 'Productivity', ACCENT,
     'Leverages robotic advantage to execute the same\nexperiment more efficiently',
     '5.05 FTE total savings realized',
     ['\u2022 Excipient Compatibility: 2.7 FTE saved',
      '\u2022 Viscous Solubility: 1.25 FTE saved',
      '\u2022 Buffer Exchange: 1.1 FTE saved']),
    ('TIER 2', 'Data & Predictive Insight', RGBColor(0x16, 0xA0, 0x85),
     'Leverages digitization advantage \u2014 larger datasets\nenable cross-project analysis and prediction',
     'Shifted EC to reactive impurity analysis',
     ['\u2022 Data aggregation across 80+ assets',
      '\u2022 Mining underlying trends',
      '\u2022 Enhanced predictive capabilities']),
    ('TIER 3', 'New Scientist-Technology Experimentation', GOLD,
     'Combines Tier\u00a01 efficiency + Tier\u00a02 insight into\nnew optimum approaches',
     'Human redundancy is NOT the goal',
     ['\u2022 New scientist\u2013technology interaction',
      '\u2022 Non-routine robotic experimentation',
      '\u2022 Automation-savvy scientist paradigm']),
]

for i, (tier, name, color, desc, highlight, bullets) in enumerate(tier_data):
    x = Inches(0.5) + i * Inches(4.1)
    y = Inches(1.6)
    w = Inches(3.9)

    # Card
    add_rounded_rect(slide, x, y, w, Inches(5.2), WHITE, color)
    # Header bar
    add_rect(slide, x, y, w, Inches(0.9), color)
    # Tier label
    add_textbox(slide, x, y + Inches(0.05), w, Inches(0.35),
                tier, font_size=13, color=RGBColor(0xFF,0xFF,0xFF),
                bold=False, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x, y + Inches(0.35), w, Inches(0.5),
                name, font_size=18, color=WHITE, bold=True,
                alignment=PP_ALIGN.CENTER)
    # Description
    add_textbox(slide, x + Inches(0.2), y + Inches(1.1), w - Inches(0.4), Inches(0.9),
                desc, font_size=13, color=DARK_TEXT, alignment=PP_ALIGN.LEFT)
    # Highlight
    add_rounded_rect(slide, x + Inches(0.15), y + Inches(2.2),
                     w - Inches(0.3), Inches(0.5), LIGHT_BG, color)
    add_textbox(slide, x + Inches(0.25), y + Inches(2.25),
                w - Inches(0.5), Inches(0.4),
                highlight, font_size=12, color=color, bold=True,
                alignment=PP_ALIGN.CENTER)
    # Bullets
    for j, bullet in enumerate(bullets):
        add_textbox(slide, x + Inches(0.2), y + Inches(2.9) + j * Inches(0.35),
                    w - Inches(0.4), Inches(0.35),
                    bullet, font_size=12, color=DARK_TEXT)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 5 — KEY DIGITAL INITIATIVES & IMPACT
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, DARK_NAVY)

add_textbox(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.6),
            'Key Digital Initiatives & Impact', font_size=36, color=WHITE, bold=True)
add_rect(slide, Inches(0.8), Inches(0.9), Inches(2.0), Inches(0.05), GOLD)

initiatives = [
    ('GPS Digital Moonshot', ACCENT,
     ['Co-led PD/CTD-wide 5-year digital vision',
      'Engaged 100+ leaders, scientists & partners',
      'Identified 21+ minimum viable products for investment',
      'Presented strategy to PDSLT and PDLT']),
    ('Bio In Silico Modeling', GREEN,
     ['Product owner \u2014 first DevEx predictive models',
      '7 dashboards, 2 production models deployed',
      '~2 FTE savings; impacted 6 major programs',
      'Set new standards for model characterization']),
    ('SPARTAN & Data Strategy', GOLD,
     ['Championed Domain-Driven Design across PD',
      'Deployed ROVER as first SPARTAN tool in PD',
      'Commissioned DPD Data Map & Maturity Framework',
      'Individual STELA Award (2024)']),
    ('ARGILE AR/VR', ORANGE,
     ['28 projects deployed across 17 global sites',
      'Critical for cross-site troubleshooting in pandemic',
      'Secured cybersecurity approvals for data capture',
      'Pioneered AR/VR in pharmaceutical development']),
    ('GenAI / LLM Strategy', RGBColor(0xAF, 0x7A, 0xC5),
     ['Collected & assessed 40+ use cases',
      'Piloted DocGenAI & ECL coding bot',
      'Launched AI literacy: office hours + newsletter',
      'Roadmap: 5 quick wins + 4 big bets']),
    ('C3.AI Data Extraction', RGBColor(0x48, 0xC9, 0xB0),
     ['Structured data from 3,000+ documents',
      '>95% accuracy on extraction',
      'Transitioned to in-house ZS platform',
      'Feeds enterprise data strategy']),
]

for i, (title, color, items) in enumerate(initiatives):
    row = i // 3
    col = i % 3
    x = Inches(0.5) + col * Inches(4.2)
    y = Inches(1.2) + row * Inches(3.0)
    w = Inches(3.9)
    h = Inches(2.7)
    # Card bg
    add_rounded_rect(slide, x, y, w, h, RGBColor(0x17, 0x30, 0x4F))
    # Left accent bar
    add_rect(slide, x, y, Inches(0.06), h, color)
    # Title
    add_textbox(slide, x + Inches(0.2), y + Inches(0.1), w - Inches(0.3), Inches(0.4),
                title, font_size=16, color=color, bold=True)
    # Bullets
    add_bullet_frame(slide, x + Inches(0.2), y + Inches(0.55), w - Inches(0.4), h - Inches(0.65),
                     items, font_size=12, color=RGBColor(0xD5, 0xDB, 0xDB),
                     bullet_char='\u25B9', spacing=3)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 6 — DIGITAL TOOLS & PLATFORMS LANDSCAPE
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, LIGHT_BG)

add_textbox(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.6),
            'Digital Tools & Platforms I\u2019ve Led or Delivered', font_size=36, color=NAVY, bold=True)
add_rect(slide, Inches(0.8), Inches(0.9), Inches(2.0), Inches(0.05), GOLD)

categories = [
    ('Data & Analytics', ACCENT,
     ['SPARTAN \u2013 PD data contextualization (DDD)',
      'ROVER \u2013 First SPARTAN production tool',
      'DevLIMS \u2013 Structured data capture',
      'Lumetics \u2013 Stability/formulation data',
      'Tabula \u2013 In-house e-batch records']),
    ('AI / ML / GenAI', GREEN,
     ['C3.AI \u2013 Document extraction',
      'Bio In Silico \u2013 Predictive models (viscosity, HMW)',
      'DocGenAI \u2013 Regulatory document drafting pilot',
      'ECL Coding Bot \u2013 GenAI for lab code',
      'NEXUS \u2013 Model launcher']),
    ('Automation Platforms', GOLD,
     ['Tecan \u2013 Liquid handling',
      'Hamilton \u2013 Liquid handling',
      'Symyx / Freeslate / Unchained Labs',
      'iAndy \u2013 Distributed automation',
      '13 Workflows \u2013 EC, viscosity, BEX, etc.']),
    ('Visualization & Operations', ORANGE,
     ['ARGILE \u2013 AR/VR (17 sites, 28 projects)',
      'Nisaba \u2013 Real-time equipment dashboards',
      'Spotfire / Tableau \u2013 Analytics',
      'Quickbase \u2013 Operational dashboards',
      'BETAR \u2013 Barriers to execution tracking']),
]

for i, (cat_name, color, tools) in enumerate(categories):
    x = Inches(0.4) + i * Inches(3.2)
    y = Inches(1.3)
    w = Inches(3.0)
    h = Inches(5.0)

    add_rounded_rect(slide, x, y, w, h, WHITE, color)
    # Header
    add_rect(slide, x, y, w, Inches(0.55), color)
    add_textbox(slide, x, y + Inches(0.08), w, Inches(0.4),
                cat_name, font_size=15, color=WHITE, bold=True,
                alignment=PP_ALIGN.CENTER)
    # Tools list
    add_bullet_frame(slide, x + Inches(0.15), y + Inches(0.7),
                     w - Inches(0.3), h - Inches(0.9),
                     tools, font_size=12, color=DARK_TEXT,
                     bullet_char='\u25AA', spacing=6)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 7 — LEADERSHIP IMPACT & RECOGNITION
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, DARK_NAVY)

add_textbox(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.6),
            'Leadership Impact & Recognition', font_size=36, color=WHITE, bold=True)
add_rect(slide, Inches(0.8), Inches(0.9), Inches(2.0), Inches(0.05), GOLD)

# Left side — People Impact
add_textbox(slide, Inches(0.8), Inches(1.3), Inches(5), Inches(0.5),
            'People & Culture', font_size=22, color=ACCENT, bold=True)

people_items = [
    'Developed Manager Readiness framework & role-competency models',
    'Grew 3 direct reports into managers in a single year',
    'Manage 7 reports (3 PhDs) + 3 people managers',
    'MyVoice engagement: 83 vs. 75 BMS avg (+8), Inclusion +11',
    'BOLD member \u2014 D&I advocate and panelist for PD leadership',
    'Created ET Culture Champion program; drove organizational dynamics',
    '13+ scientists rotated through automation \u2192 3 publications',
]
add_bullet_frame(slide, Inches(0.8), Inches(1.9), Inches(5.5), Inches(3.5),
                 people_items, font_size=14, color=RGBColor(0xD5,0xDB,0xDB),
                 bullet_char='\u25B8', spacing=6)

# Right side — Recognition
add_textbox(slide, Inches(7.0), Inches(1.3), Inches(5.5), Inches(0.5),
            'Recognition', font_size=22, color=GOLD, bold=True)

recognition = [
    'STELA Award (Individual, 2024) \u2014 SPARTAN adoption',
    '\u201cPhenomenal performance\u201d \u2014 Manager Reviews (2024, 2025)',
    '\u201cInspiring and visionary leader\u201d \u2014 Promotion Citation',
    '\u201cCreated Digital Engineering in DPST\u201d \u2014 Senior Leadership',
    'AAPS webinar presenter & automation user-group speaker',
    'Vendor co-development \u2192 commercially deployed platform',
    '3 promotions in 5 years (SRS II \u2192 Principal \u2192 Sci. Director)',
]
add_bullet_frame(slide, Inches(7.0), Inches(1.9), Inches(5.8), Inches(3.5),
                 recognition, font_size=14, color=RGBColor(0xD5,0xDB,0xDB),
                 bullet_char='\u25B8', spacing=6)

# Bottom quote box
add_rounded_rect(slide, Inches(0.8), Inches(5.7), Inches(11.7), Inches(1.2),
                 RGBColor(0x17, 0x30, 0x4F), GOLD)
add_textbox(slide, Inches(1.0), Inches(5.85), Inches(11.3), Inches(0.9),
            '\u201cHe has a vision of a future for DPST, PD, BMS and our products '
            'and he truly wants to shape the future of healthcare.\u201d\n'
            '\u2014 Martin McLoughlin, Head of ET & Device Development',
            font_size=16, color=WHITE, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 8 — WHAT I'M LOOKING FOR IN A MENTOR
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, LIGHT_BG)

add_textbox(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.6),
            'What I\u2019m Looking For', font_size=36, color=NAVY, bold=True)
add_rect(slide, Inches(0.8), Inches(0.9), Inches(2.0), Inches(0.05), GOLD)
add_textbox(slide, Inches(0.8), Inches(1.05), Inches(10), Inches(0.4),
            'Perspectives and guidance as I move toward executive digital leadership in pharma',
            font_size=15, color=MED_GRAY)

topics = [
    ('Scaling AI in R&D', ACCENT,
     'Moving from pilots and proofs-of-concept to enterprise-wide AI adoption \u2014 '
     'governance, change management, and measuring ROI at scale.'),
    ('Executive Presence', GREEN,
     'Strengthening my voice at the VP+ level \u2014 strategic framing, board-level '
     'communication, and influencing without direct authority across matrixed organizations.'),
    ('Building the Bridge', GOLD,
     'Best practices for sustaining productive IT\u2013Science partnerships \u2014 '
     'balancing enterprise architecture standards with scientific speed-to-value.'),
    ('Career Trajectory', ORANGE,
     'Navigating the path from Scientific Director to VP-level digital/AI leadership roles \u2014 '
     'visibility, sponsorship, and portfolio expansion.'),
]

for i, (topic, color, description) in enumerate(topics):
    row = i // 2
    col = i % 2
    x = Inches(0.5) + col * Inches(6.2)
    y = Inches(1.7) + row * Inches(2.6)
    w = Inches(5.9)
    h = Inches(2.3)

    add_rounded_rect(slide, x, y, w, h, WHITE, color)
    # Icon circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                     x + Inches(0.25), y + Inches(0.3),
                                     Inches(0.5), Inches(0.5))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    # Number in circle
    add_textbox(slide, x + Inches(0.25), y + Inches(0.33),
                Inches(0.5), Inches(0.45),
                str(i + 1), font_size=20, color=WHITE, bold=True,
                alignment=PP_ALIGN.CENTER)
    # Topic title
    add_textbox(slide, x + Inches(0.95), y + Inches(0.3), w - Inches(1.2), Inches(0.4),
                topic, font_size=18, color=color, bold=True)
    # Description
    add_textbox(slide, x + Inches(0.95), y + Inches(0.8), w - Inches(1.2), Inches(1.3),
                description, font_size=13, color=DARK_TEXT)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 9 — CLOSING / CONTACT
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, DARK_NAVY)

add_rect(slide, Inches(4.5), Inches(1.8), Inches(4.3), Inches(0.05), GOLD)

add_textbox(slide, Inches(0.5), Inches(2.0), Inches(12.3), Inches(1.0),
            'Thank You', font_size=52, color=WHITE, bold=True,
            alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(0.5), Inches(3.2), Inches(12.3), Inches(0.6),
            'I\u2019m looking forward to learning from your experience.',
            font_size=22, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Contact info
contact_lines = [
    'Jaquan K. Levons',
    '609.240.1241  |  jaquan.levons@gmail.com',
    'linkedin.com/in/jaquan-levons',
]
for i, line in enumerate(contact_lines):
    fs = 20 if i == 0 else 16
    clr = WHITE if i == 0 else ACCENT
    bld = True if i == 0 else False
    add_textbox(slide, Inches(0.5), Inches(4.5) + i * Inches(0.45),
                Inches(12.3), Inches(0.4), line, font_size=fs, color=clr,
                bold=bld, alignment=PP_ALIGN.CENTER)


# ── Save ────────────────────────────────────────────────────────────────────
out_dir = r'C:\Users\jaqua\Downloads\JKL_Resumes'
out_path = os.path.join(out_dir, 'Levons_Digital_Experience_Mentor_MeetGreet.pptx')
prs.save(out_path)
print('Saved to:', out_path)
